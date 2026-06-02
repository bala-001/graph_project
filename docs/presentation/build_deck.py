"""Build the PAIQ Idea 11.3 Phase 1 presentation (python-pptx).

Run:  python docs/presentation/build_deck.py
Out:  docs/presentation/PAIQ-11.3-Phase1.pptx

Content is grounded in the actual code (src/d_extraction, src/guardrails,
src/journal, src/shadow) and docs/ (design.md, architecture/). No em dashes.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette -------------------------------------------------------------
INK = RGBColor(0x14, 0x32, 0x3B)        # deep teal-navy (primary dark)
TEAL = RGBColor(0x27, 0x78, 0x84)       # mid teal
TEAL_LT = RGBColor(0x5E, 0xA8, 0xA7)    # light teal
CORAL = RGBColor(0xFE, 0x44, 0x47)      # accent: catches / rejects / kill gates
AMBER = RGBColor(0xF2, 0xA3, 0x3C)      # secondary accent: retry / warn
BG = RGBColor(0xF4, 0xF6, 0xF6)         # off-white slide background
CARD = RGBColor(0xFF, 0xFF, 0xFF)       # card fill
GRAY = RGBColor(0x6B, 0x7B, 0x83)       # muted text
LINE = RGBColor(0xD4, 0xDD, 0xDE)       # hairlines
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = "Arial"
BODY = "Arial"
MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---- primitives ----------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _set_text(tf, lines, size, color, bold=False, font=BODY, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.0):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing:
            p.line_spacing = line_spacing
        # support (text, opts) tuples per-line
        runs = ln if isinstance(ln, list) else [(ln, {})]
        for text, opts in runs:
            r = p.add_run(); r.text = text
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.name = opts.get("font", font)
            r.font.color.rgb = opts.get("color", color)
    return tf


def textbox(s, x, y, w, h, lines, size=18, color=INK, bold=False, font=BODY,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    _set_text(tb.text_frame, lines, size, color, bold, font, align, anchor,
              space_after, line_spacing)
    return tb


def rect(s, x, y, w, h, fill, line=None, line_w=0, rounded=False, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w or 1)
    shp.shadow.inherit = False
    return shp


def card(s, x, y, w, h, fill=CARD, bar=None):
    rect(s, x, y, w, h, fill, line=LINE, line_w=1, rounded=True)
    if bar:
        rect(s, x, y, Inches(0.09), h, bar, rounded=False)


def chevron(s, x, y, w, h, color):
    a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    return a


def header(s, kicker, title, n):
    rect(s, 0, 0, Inches(0.22), EMU_H, TEAL)          # left spine
    textbox(s, Inches(0.55), Inches(0.34), Inches(10), Inches(0.34),
            kicker.upper(), size=12, color=TEAL, bold=True)
    textbox(s, Inches(0.5), Inches(0.6), Inches(11.6), Inches(0.95),
            title, size=30, color=INK, bold=True, font=HEAD, line_spacing=0.98)
    rect(s, Inches(0.55), Inches(1.5), Inches(1.1), Pt(3), CORAL)
    textbox(s, Inches(12.2), Inches(0.36), Inches(0.9), Inches(0.4),
            f"{n:02d}", size=13, color=GRAY, bold=True, align=PP_ALIGN.RIGHT)


def bullets(s, x, y, w, items, size=16, gap=8):
    tb = s.shapes.add_textbox(x, y, w, Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lead, rest = (it if isinstance(it, tuple) else (it, ""))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.05
        r = p.add_run(); r.text = "— "
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = CORAL; r.font.name = BODY
        r2 = p.add_run(); r2.text = lead
        r2.font.size = Pt(size); r2.font.bold = bool(rest); r2.font.color.rgb = INK; r2.font.name = BODY
        if rest:
            r3 = p.add_run(); r3.text = "  " + rest
            r3.font.size = Pt(size); r3.font.color.rgb = GRAY; r3.font.name = BODY
    return tb


# =========================================================================
# 1. TITLE
# =========================================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, INK)
rect(s, 0, Inches(5.0), EMU_W, Inches(2.5), RGBColor(0x0F, 0x27, 0x2E))
rect(s, Inches(0.9), Inches(1.7), Inches(1.5), Pt(5), CORAL)
textbox(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.5),
        "PAIQ IDEA 11.3  ·  PHASE 1", size=15, color=TEAL_LT, bold=True)
textbox(s, Inches(0.85), Inches(2.45), Inches(11.7), Inches(2.0),
        ["Relationship-Aware Extraction", "+ Extraction-Time Guardrails"],
        size=42, color=WHITE, bold=True, font=HEAD, line_spacing=1.0)
textbox(s, Inches(0.9), Inches(4.35), Inches(11.2), Inches(0.7),
        "Emit structured edge-triples natively, and catch logical "
        "contradictions before they are committed.",
        size=18, color=TEAL_LT)
textbox(s, Inches(0.9), Inches(5.45), Inches(11.5), Inches(1.4),
        [
            [("Approach D", {"bold": True, "color": WHITE, "size": 15}),
             ("   modified extraction emits fields + edges via provider structured outputs", {"color": RGBColor(0xB8, 0xCB, 0xCE), "size": 15})],
            [("Guardrails", {"bold": True, "color": WHITE, "size": 15}),
             ("   4 detectors fire after each LLM call, before the next chunk commits", {"color": RGBColor(0xB8, 0xCB, 0xCE), "size": 15})],
        ], size=15, space_after=8)
textbox(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
        "github.com/bala-001/graph_project   ·   Phase 1 buildable-now subset   ·   Status: unit-green",
        size=11, color=GRAY)

# =========================================================================
# 2. THE PROBLEM
# =========================================================================
s = slide()
header(s, "Why this exists", "Clients see contradictory and incomplete criteria", 2)
bullets(s, Inches(0.55), Inches(1.85), Inches(6.3), [
    ("Facts are extracted correctly today.", "Drugs, age limits, dates, step-therapy entries, quantity limits."),
    ("Relationships between them are lost.", "Prerequisite chains split into disconnected facts."),
    ("Contradictions ship unflagged.", "Two criteria for the same drug-indication conflict, no warning."),
    ("Multi-page content is handled chunk-by-chunk.", "A definition on page 3 modifying a criterion on page 17 is not reconstructed."),
    ("External pressure.", "Client complaint tickets are the demand signal, not internal opinion."),
], size=15, gap=11)
# callout card
card(s, Inches(7.2), Inches(1.95), Inches(5.6), Inches(4.4), bar=CORAL)
textbox(s, Inches(7.5), Inches(2.15), Inches(5.1), Inches(0.4),
        "EXAMPLE: SPLIT PREREQUISITE CHAIN", size=12, color=CORAL, bold=True)
textbox(s, Inches(7.5), Inches(2.6), Inches(5.1), Inches(2.0),
        [
            [("Drug A", {"font": MONO, "bold": True, "color": INK, "size": 15}),
             (" requires step therapy with ", {"color": GRAY, "size": 15}),
             ("Drug B", {"font": MONO, "bold": True, "color": INK, "size": 15})],
            [("Drug B", {"font": MONO, "bold": True, "color": INK, "size": 15}),
             (" has an age restriction", {"color": GRAY, "size": 15})],
        ], size=15, space_after=10, line_spacing=1.1)
rect(s, Inches(7.5), Inches(3.7), Inches(5.0), Pt(1), LINE)
textbox(s, Inches(7.5), Inches(3.85), Inches(5.1), Inches(2.3),
        [
            [("Today: ", {"bold": True, "color": INK, "size": 14}),
             ("3 disconnected facts. The chain A → B → age is never surfaced.", {"color": GRAY, "size": 14})],
            [("Goal: ", {"bold": True, "color": TEAL, "size": 14}),
             ("emit the edges, detect the contradictions, reduce complaint tickets by 50% in 90 days.", {"color": INK, "size": 14})],
        ], size=14, space_after=12, line_spacing=1.12)

# =========================================================================
# 3. SOLUTION AT A GLANCE
# =========================================================================
s = slide()
header(s, "The approach", "Two changes to the extraction pipeline", 3)
# pillar A
card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.5), bar=TEAL)
textbox(s, Inches(0.9), Inches(2.15), Inches(5.4), Inches(0.5), "APPROACH D", size=13, color=TEAL, bold=True)
textbox(s, Inches(0.9), Inches(2.55), Inches(5.4), Inches(0.6), "Relationship-aware extraction", size=20, color=INK, bold=True)
bullets(s, Inches(0.9), Inches(3.3), Inches(5.3), [
    ("Emits fields + edges natively.", "Same call, structured edge-triples added."),
    ("Provider built-in structured outputs.", "OpenAI json_schema strict OR Anthropic tool-use."),
    ("Pydantic-defined schema.", "5 predicates, typed nodes, qualifiers."),
    ("Behind a feature flag.", "paiq.d_extraction.enabled = instant revert."),
], size=14, gap=10)
# pillar B
card(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(4.5), bar=CORAL)
textbox(s, Inches(7.15), Inches(2.15), Inches(5.4), Inches(0.5), "GUARDRAILS", size=13, color=CORAL, bold=True)
textbox(s, Inches(7.15), Inches(2.55), Inches(5.4), Inches(0.6), "Extraction-time consistency check", size=20, color=INK, bold=True)
bullets(s, Inches(7.15), Inches(3.3), Inches(5.3), [
    ("4 detection scenarios.", "Circular, prereq mismatch, contradictory limits, age conflict."),
    ("Fires between LLM calls.", "After each chunk, before the next commits (D10)."),
    ("Reject + retry with validator error.", "Up to 3 retries, then analyst flag."),
    ("3-counter FP telemetry.", "Clean false-positive rate for the Week-6 gate."),
], size=14, gap=10)

# =========================================================================
# 4. EDGE DATA MODEL
# =========================================================================
s = slide()
header(s, "Data model", "What an edge looks like", 4)
# left: predicates
textbox(s, Inches(0.55), Inches(1.85), Inches(5.2), Inches(0.4),
        "FIVE EDGE PREDICATES", size=13, color=TEAL, bold=True)
preds = [("requires", "A needs B satisfied"), ("excludes", "A rules out B"),
         ("applies_to", "criterion scope"), ("overrides", "amendment supersedes"),
         ("effective_from", "date a rule starts")]
yy = Inches(2.3)
for name, desc in preds:
    card(s, Inches(0.55), yy, Inches(5.6), Inches(0.66), bar=TEAL_LT)
    textbox(s, Inches(0.85), yy + Inches(0.06), Inches(2.2), Inches(0.55),
            name, size=16, color=INK, bold=True, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(3.0), yy + Inches(0.06), Inches(3.0), Inches(0.55),
            desc, size=13, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
    yy += Inches(0.8)
# right: schema card
card(s, Inches(6.5), Inches(1.95), Inches(6.3), Inches(4.6), fill=INK)
textbox(s, Inches(6.85), Inches(2.15), Inches(5.7), Inches(0.4),
        "Edge  (src/d_extraction/schema.py)", size=13, color=TEAL_LT, bold=True, font=MONO)
code = [
    "kind:        EdgeKind        # one of 5",
    "subject:     DrugNode        # canonical_id",
    "object:      DrugNode |",
    "             IndicationNode | str",
    "qualifiers:  dict            # age_min/max,",
    "             dosage_*, quantity_*, age_exact",
    "source_page / source_chunk_id",
    "model_confidence: float|None",
]
textbox(s, Inches(6.85), Inches(2.65), Inches(5.7), Inches(2.6),
        code, size=13, color=WHITE, font=MONO, space_after=5, line_spacing=1.05)
rect(s, Inches(6.85), Inches(5.15), Inches(5.6), Pt(1), TEAL)
textbox(s, Inches(6.85), Inches(5.3), Inches(5.7), Inches(1.1),
        [
            [("Shape, not semantics. ", {"bold": True, "color": AMBER, "size": 13})],
            [("Strict schema guarantees a valid edge can still be logically wrong. Guardrails + eval are the semantic defense.", {"color": RGBColor(0xCF, 0xDC, 0xDE), "size": 13})],
        ], size=13, space_after=4, line_spacing=1.1)

# =========================================================================
# 5. END-TO-END WORKFLOW
# =========================================================================
s = slide()
header(s, "How it runs", "End-to-end pipeline", 5)
textbox(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.4),
        "Per document: chunk → extract → guard → journal → materialize. Edges are checked before the next chunk commits.",
        size=14, color=GRAY)
# pipeline boxes
steps = [
    ("PBM\ndocument", TEAL),
    ("Chunk\ni..N", TEAL),
    ("D-mode\nLLM call", TEAL),
    ("Guardrails\ncheck", CORAL),
    ("Batched\njournal", TEAL),
    ("Materialize\n+ complete", INK),
]
x = Inches(0.55); y = Inches(2.55); bw = Inches(1.74); bh = Inches(1.2); gap = Inches(0.27)
for i, (label, col) in enumerate(steps):
    card(s, x, y, bw, bh, fill=WHITE, bar=col)
    textbox(s, x, y, bw, bh, label.split("\n"), size=14, color=INK, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
    if i < len(steps) - 1:
        chevron(s, x + bw + Inches(0.02), y + Inches(0.42), gap - Inches(0.04), Inches(0.36), TEAL_LT)
    x += bw + gap
# loop annotation
textbox(s, Inches(2.9), Inches(3.85), Inches(5.6), Inches(0.4),
        "↺  retry / loop per chunk (max 3)", size=12, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
# downstream row
card(s, Inches(0.55), Inches(4.55), Inches(6.0), Inches(1.85), bar=TEAL)
textbox(s, Inches(0.85), Inches(4.7), Inches(5.4), Inches(0.4), "DOWNSTREAM CONSUMERS", size=12, color=TEAL, bold=True)
textbox(s, Inches(0.85), Inches(5.1), Inches(5.5), Inches(1.2),
        [
            [("Check ", {"color": INK, "size": 14}), ("extraction_complete", {"font": MONO, "bold": True, "color": INK, "size": 14}),
             (" first.", {"color": INK, "size": 14})],
            [("false → row is INVISIBLE (not-yet-extracted), 24h GC re-queues.", {"color": GRAY, "size": 13})],
            [("true → surface fields + edges to analyst tools / client reports.", {"color": GRAY, "size": 13})],
        ], size=14, space_after=7, line_spacing=1.1)
card(s, Inches(6.8), Inches(4.55), Inches(6.0), Inches(1.85), bar=CORAL)
textbox(s, Inches(7.1), Inches(4.7), Inches(5.4), Inches(0.4), "FEATURE FLAG", size=12, color=CORAL, bold=True)
textbox(s, Inches(7.1), Inches(5.1), Inches(5.5), Inches(1.2),
        [
            [("paiq.d_extraction.enabled", {"font": MONO, "bold": True, "color": INK, "size": 14})],
            [("true  → D prompts run, edges emitted, guardrails consulted.", {"color": GRAY, "size": 13})],
            [("false → baseline prompts, no edges. Instant revert path.", {"color": GRAY, "size": 13})],
        ], size=14, space_after=7, line_spacing=1.1)

# =========================================================================
# 6. THE GUARDRAILS LOOP
# =========================================================================
s = slide()
header(s, "Guardrails", "Detection order, then verdict", 6)
textbox(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.4),
        "Four scenarios run in a fixed order. The first hit short-circuits to REJECT_RETRY. No hit means ACCEPT.",
        size=14, color=GRAY)
det = [("1  Circular dependency", "A requires B; B requires A"),
       ("2  Prerequisite mismatch", "requires vs excludes on same pair"),
       ("3  Contradictory limits", "min bound exceeds max bound"),
       ("4  Age conflict", "exact age qualifier mismatch")]
x = Inches(0.55); y = Inches(2.45); bw = Inches(2.92); bh = Inches(1.35); gap = Inches(0.13)
for i, (t, d) in enumerate(det):
    card(s, x, y, bw, bh, bar=CORAL)
    textbox(s, x + Inches(0.18), y + Inches(0.16), bw - Inches(0.3), Inches(0.5), t, size=14, color=INK, bold=True)
    textbox(s, x + Inches(0.18), y + Inches(0.62), bw - Inches(0.3), Inches(0.7), d, size=12, color=GRAY, line_spacing=1.05)
    if i < 3:
        chevron(s, x + bw + Inches(0.005), y + Inches(0.5), gap, Inches(0.34), TEAL_LT)
    x += bw + gap
# verdict row
card(s, Inches(0.55), Inches(4.2), Inches(6.0), Inches(2.2), fill=RGBColor(0xEC, 0xF4, 0xF3), bar=TEAL)
textbox(s, Inches(0.85), Inches(4.4), Inches(5.4), Inches(0.5), "ACCEPT", size=18, color=TEAL, bold=True)
textbox(s, Inches(0.85), Inches(4.95), Inches(5.5), Inches(1.3),
        "No detector triggered. The edge is appended to the journal buffer and "
        "extraction proceeds to the next chunk.", size=14, color=INK, line_spacing=1.15)
card(s, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.2), fill=RGBColor(0xFD, 0xEC, 0xEC), bar=CORAL)
textbox(s, Inches(7.1), Inches(4.4), Inches(5.4), Inches(0.5), "REJECT_RETRY", size=18, color=CORAL, bold=True)
textbox(s, Inches(7.1), Inches(4.95), Inches(5.5), Inches(1.3),
        "Re-extract this relationship with the validator error attached. "
        "Up to 3 retries, then fall back to an analyst flag.", size=14, color=INK, line_spacing=1.15)

# =========================================================================
# 7. RETRY + 3-COUNTER FP TELEMETRY
# =========================================================================
s = slide()
header(s, "Telemetry", "Three outcomes, one clean false-positive rate", 7)
textbox(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.4),
        "On REJECT_RETRY the model re-extracts. Each rejected edge resolves to exactly one counter (D4 + canonicalization).",
        size=14, color=GRAY)
cols = [
    ("COUNTER 1", "Same edge again", "Model reproduces the same canonical edge. The guardrail was likely wrong. Real false positive. Accept the edge.", CORAL),
    ("COUNTER 2", "Different, now passes", "Model returns a corrected edge that passes the check. True positive. Accept it.", TEAL),
    ("COUNTER 3", "Retries exhausted", "No acceptable edge after MAX_RETRIES = 3. Raise an analyst flag. Accept nothing.", AMBER),
]
x = Inches(0.55); w = Inches(3.97); gap = Inches(0.2)
for title, sub, body, col in cols:
    card(s, x, Inches(2.4), w, Inches(2.7), bar=col)
    textbox(s, x + Inches(0.25), Inches(2.6), w - Inches(0.5), Inches(0.4), title, size=13, color=col, bold=True)
    textbox(s, x + Inches(0.25), Inches(3.0), w - Inches(0.5), Inches(0.5), sub, size=17, color=INK, bold=True)
    textbox(s, x + Inches(0.25), Inches(3.6), w - Inches(0.5), Inches(1.4), body, size=13, color=GRAY, line_spacing=1.15)
    x += w + gap
# formula band
rect(s, Inches(0.55), Inches(5.4), Inches(12.25), Inches(1.0), INK, rounded=True)
textbox(s, Inches(0.9), Inches(5.55), Inches(11.6), Inches(0.75),
        [
            [("FP rate  =  counter 1 / total", {"font": MONO, "bold": True, "color": WHITE, "size": 18}),
             ("        Week-6 Kill Criteria gate: must stay within ", {"color": RGBColor(0xCF, 0xDC, 0xDE), "size": 15}),
             ("1–15%", {"bold": True, "color": AMBER, "size": 16}),
             (".  Outside → tune or disable.", {"color": RGBColor(0xCF, 0xDC, 0xDE), "size": 15})],
        ], size=16, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================================
# 8. PERSIST-AS-YOU-GO / JOURNAL
# =========================================================================
s = slide()
header(s, "Durability", "Persist-as-you-go journal", 8)
bullets(s, Inches(0.55), Inches(1.9), Inches(6.2), [
    ("Accepted edges buffer in memory.", "One JournalWriter per document."),
    ("Flush every 10 edges OR 5 seconds.", "Whichever comes first (D7 batching)."),
    ("One JSON line per edge.", "Appended to document_id.journal."),
    ("On complete: replay + collapse.", "Write document_id.json, set complete = true."),
    ("Crash mid-document is safe.", "Flag stays false; row invisible; 24h GC re-queues."),
], size=15, gap=12)
# flow on the right
rx = Inches(7.1); rw = Inches(5.5)
def fbox(y, label, col, sub=None):
    h = Inches(0.72) if sub else Inches(0.6)
    card(s, rx, y, rw, h, bar=col)
    textbox(s, rx + Inches(0.25), y + Inches(0.07), rw - Inches(0.5), Inches(0.4),
            label, size=14, color=INK, bold=True)
    if sub:
        textbox(s, rx + Inches(0.25), y + Inches(0.4), rw - Inches(0.5), Inches(0.3),
                sub, size=11, color=GRAY)
    return y + h
y = Inches(1.95)
y = fbox(y, "emit edge → in-memory buffer", TEAL); y += Inches(0.18)
y = fbox(y, "flush when  len ≥ 10  OR  5s elapsed", AMBER, "append batch to .journal"); y += Inches(0.18)
y = fbox(y, "extraction_complete = true", INK, "final flush → replay → .json"); y += Inches(0.18)
y = fbox(y, "downstream sees the document", TEAL, "only after .json exists")

# =========================================================================
# 9. SHADOW MODE + MEASUREMENT
# =========================================================================
s = slide()
header(s, "Validation", "Shadow mode: measure D before it ships", 9)
textbox(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.4),
        "Baseline and D-mode run on the same document. The harness compares against analyst-corrected ground truth.",
        size=14, color=GRAY)
# two comparison cards
card(s, Inches(0.55), Inches(2.35), Inches(3.9), Inches(2.0), bar=GRAY)
textbox(s, Inches(0.8), Inches(2.5), Inches(3.4), Inches(0.4), "BASELINE", size=13, color=GRAY, bold=True)
textbox(s, Inches(0.8), Inches(2.9), Inches(3.4), Inches(1.3), "Existing prompts. Fields only, no edges. The regression reference point.", size=14, color=INK, line_spacing=1.15)
card(s, Inches(4.65), Inches(2.35), Inches(3.9), Inches(2.0), bar=TEAL)
textbox(s, Inches(4.9), Inches(2.5), Inches(3.4), Inches(0.4), "D-MODE", size=13, color=TEAL, bold=True)
textbox(s, Inches(4.9), Inches(2.9), Inches(3.4), Inches(1.3), "Fields + edges. Compared edge-by-edge using canonicalized membership.", size=14, color=INK, line_spacing=1.15)
card(s, Inches(8.75), Inches(2.35), Inches(4.05), Inches(2.0), bar=CORAL)
textbox(s, Inches(9.0), Inches(2.5), Inches(3.5), Inches(0.4), "GROUND TRUTH", size=13, color=CORAL, bold=True)
textbox(s, Inches(9.0), Inches(2.9), Inches(3.5), Inches(1.3), "Analyst-corrected edges on the 50-doc cascade-OCR eval set. Live run is Phase-0-gated.", size=14, color=INK, line_spacing=1.15)
# metrics band
mets = [
    ("Edge precision / recall", "|D ∩ GT| / |D|  and  / |GT|"),
    ("Wilson 95% lower bound", "defends gates from small-N noise"),
    ("Field iso-precision", "regression if a field value changes > 5%"),
]
x = Inches(0.55); w = Inches(4.03)
for t, d in mets:
    card(s, x, Inches(4.6), w, Inches(1.8), fill=RGBColor(0xEC, 0xF4, 0xF3))
    textbox(s, x + Inches(0.25), Inches(4.8), w - Inches(0.5), Inches(0.8), t, size=16, color=INK, bold=True, line_spacing=1.0)
    textbox(s, x + Inches(0.25), Inches(5.55), w - Inches(0.5), Inches(0.7), d, size=13, color=TEAL, font=MONO, line_spacing=1.05)
    x += w + Inches(0.18)

# =========================================================================
# 10. MODULE ARCHITECTURE
# =========================================================================
s = slide()
header(s, "Architecture", "Module boundaries", 10)
textbox(s, Inches(0.55), Inches(1.7), Inches(12), Inches(0.4),
        "Clear lanes so two engineers can work in parallel. Dependencies point one way.",
        size=14, color=GRAY)
mods = [
    ("src/d_extraction", "Edge schema, D-mode + baseline prompts, multi-call loop", TEAL),
    ("src/guardrails", "4 detectors, retry policy, 3-counter FP telemetry, state", CORAL),
    ("src/journal", "Batched-write journal + crash-recovery replay", TEAL),
    ("src/shadow", "D-vs-ground-truth comparison + Wilson bounds", TEAL_LT),
    ("src/feature_flags", "paiq.d_extraction.enabled gates prompt-mode swap", AMBER),
    ("src/telemetry", "FP-rate counters, observability hooks", TEAL_LT),
    ("src/cascade_integration", "Cascade-OCR judge re-calibration (one-time, pre-ship)", GRAY),
]
x = Inches(0.55); y = Inches(2.4); w = Inches(4.03); h = Inches(1.25)
gx = Inches(0.18); gy = Inches(0.2)
for i, (name, desc, col) in enumerate(mods):
    cx = x + (i % 3) * (w + gx)
    cy = y + (i // 3) * (h + gy)
    card(s, cx, cy, w, h, bar=col)
    textbox(s, cx + Inches(0.22), cy + Inches(0.16), w - Inches(0.4), Inches(0.4),
            name, size=15, color=INK, bold=True, font=MONO)
    textbox(s, cx + Inches(0.22), cy + Inches(0.6), w - Inches(0.4), Inches(0.6),
            desc, size=12.5, color=GRAY, line_spacing=1.08)

# =========================================================================
# 11. KILL CRITERIA
# =========================================================================
s = slide()
header(s, "Risk gates", "Kill Criteria keep the bet survivable", 11)
textbox(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.4),
        "Each gate has a trigger and a pre-agreed decision. Fire one, and the project pauses or reverts.",
        size=14, color=GRAY)
rows = [
    ("Week 4", "Field precision regresses > 5% on existing types", "Flip the flag off. Instant revert.", CORAL),
    ("Week 6", "Guardrail FP rate outside 1–15% (counter 1 / total)", "Tune the detectors or disable.", AMBER),
    ("Quarter 1", "Edge precision < 85% OR recall < 80% on shadow", "No Phase 2 expansion.", CORAL),
    ("Week 1", "No written sponsor commitment in 7 days", "Downgrade to contradiction-only wedge.", GRAY),
]
y = Inches(2.45); rw = Inches(12.25); rh = Inches(0.92)
# header strip
rect(s, Inches(0.55), y, rw, Inches(0.5), INK, rounded=False)
textbox(s, Inches(0.75), y + Inches(0.06), Inches(1.6), Inches(0.4), "GATE", size=12, color=WHITE, bold=True)
textbox(s, Inches(2.5), y + Inches(0.06), Inches(6.0), Inches(0.4), "TRIGGER", size=12, color=WHITE, bold=True)
textbox(s, Inches(9.2), y + Inches(0.06), Inches(3.4), Inches(0.4), "DECISION", size=12, color=WHITE, bold=True)
y += Inches(0.5)
for gate, trig, dec, col in rows:
    card(s, Inches(0.55), y, rw, rh)
    rect(s, Inches(0.55), y, Inches(0.09), rh, col)
    textbox(s, Inches(0.75), y, Inches(1.7), rh, gate, size=15, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(2.5), y, Inches(6.5), rh, trig, size=13.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    textbox(s, Inches(9.2), y, Inches(3.5), rh, dec, size=13.5, color=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    y += rh + Inches(0.13)

# =========================================================================
# 12. STATUS + NEXT STEPS
# =========================================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, INK)
rect(s, Inches(0.9), Inches(0.9), Inches(1.5), Pt(5), CORAL)
textbox(s, Inches(0.9), Inches(1.15), Inches(11), Inches(0.5), "WHERE WE ARE", size=14, color=TEAL_LT, bold=True)
textbox(s, Inches(0.85), Inches(1.6), Inches(11.6), Inches(0.9), "Status and next steps", size=32, color=WHITE, bold=True, font=HEAD)
# done card
card(s, Inches(0.9), Inches(2.75), Inches(5.7), Inches(3.4), fill=RGBColor(0x0F, 0x27, 0x2E), bar=TEAL)
textbox(s, Inches(1.2), Inches(2.95), Inches(5.1), Inches(0.4), "DONE  ·  BUILDABLE-NOW SUBSET", size=13, color=TEAL_LT, bold=True)
textbox(s, Inches(1.2), Inches(3.45), Inches(5.2), Inches(2.6),
        [
            [("✓  Edge schema + canonicalization", {"color": WHITE, "size": 15})],
            [("✓  4 guardrail detectors", {"color": WHITE, "size": 15})],
            [("✓  Retry policy + 3-counter telemetry", {"color": WHITE, "size": 15})],
            [("✓  Batched journal + replay", {"color": WHITE, "size": 15})],
            [("✓  Shadow comparison + Wilson bounds", {"color": WHITE, "size": 15})],
            [("✓  Feature flag + downstream visibility", {"color": WHITE, "size": 15})],
            [("✓  Unit-green test suite", {"color": WHITE, "size": 15})],
        ], size=15, space_after=7)
# next card
card(s, Inches(6.9), Inches(2.75), Inches(5.5), Inches(3.4), fill=RGBColor(0x0F, 0x27, 0x2E), bar=CORAL)
textbox(s, Inches(7.2), Inches(2.95), Inches(4.9), Inches(0.4), "NEXT  ·  PHASE-0 GATED", size=13, color=AMBER, bold=True)
textbox(s, Inches(7.2), Inches(3.45), Inches(5.0), Inches(2.6),
        [
            [("→  Close Q3: enumerate D-mode prompts/schemas", {"color": WHITE, "size": 15})],
            [("→  Wire CI eval gate (GitHub Actions)", {"color": WHITE, "size": 15})],
            [("→  Cascade-OCR judge re-calibration run", {"color": WHITE, "size": 15})],
            [("→  Live shadow run on 50-doc eval set", {"color": WHITE, "size": 15})],
            [("→  Complaint root-cause audit (Week-2 gate)", {"color": WHITE, "size": 15})],
            [("→  Sponsor written commitment (Week-1 gate)", {"color": WHITE, "size": 15})],
        ], size=15, space_after=8)
textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6),
        "Primary success metric: reduce contradiction / missed-relationship complaint tickets by 50% within 90 days of shadow launch.",
        size=13, color=TEAL_LT)

# ---- save ----------------------------------------------------------------
out = Path(__file__).resolve().parent / "PAIQ-11.3-Phase1.pptx"
prs.save(str(out))
print(f"saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
